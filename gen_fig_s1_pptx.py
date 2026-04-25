#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generate Supplementary Figure S1 as editable PPTX.
Top: Three reasoning chains (Published Evidence → Consequence → Weight Adjustment)
Bottom: Patient A vs Patient B worked example

Data provenance: Table 3 weights; [11] Aubrey 2018, [20] Solomon 2011, [9] StatPearls, [21] Lee 2004
"""
import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

TEAL=RGBColor(0x1A,0x7A,0x5C); CORAL=RGBColor(0xE0,0x7B,0x54)
PURPLE=RGBColor(0x7B,0x5E,0xA7); NAVY=RGBColor(0x2C,0x3E,0x6B)
LTEAL=RGBColor(0xD0,0xED,0xE3); LCORAL=RGBColor(0xFD,0xEB,0xD8)
LPURP=RGBColor(0xE8,0xDF,0xF0)
BODY=RGBColor(0x33,0x33,0x33); GRAY=RGBColor(0x66,0x66,0x66)
WHITE=RGBColor(0xFF,0xFF,0xFF); GRAYBG=RGBColor(0xF8,0xF8,0xF8)
FONT='Calibri'

def txt(sl,x,y,w,h,text,size=8,bold=False,italic=False,color=BODY,align=PP_ALIGN.CENTER):
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.name=FONT; r.font.size=Pt(size)
    r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color
    return tb

def rrect(sl,x,y,w,h,fill,line,lw=1.5):
    s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb=fill; s.line.color.rgb=line; s.line.width=Pt(lw)
    s.adjustments[0]=0.08; return s

def arrow_r(sl,x1,y,x2,color):
    """Right-pointing arrow connector."""
    c=sl.shapes.add_connector(1,Inches(x1),Inches(y),Inches(x2),Inches(y))
    c.line.color.rgb=color; c.line.width=Pt(1.5)
    # Add arrowhead
    s=sl.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,Inches(x2-0.08),Inches(y-0.06),Inches(0.12),Inches(0.12))
    s.rotation=90.0; s.fill.solid(); s.fill.fore_color.rgb=color; s.line.fill.background()

def ln(sl,x1,y1,x2,y2,color,width=1.0):
    c=sl.shapes.add_connector(1,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    c.line.color.rgb=color; c.line.width=Pt(width)

prs=Presentation()
prs.slide_width=Inches(10)
prs.slide_height=Inches(12)
sl=prs.slides.add_slide(prs.slide_layouts[6])

# ── Panel A title ──
txt(sl,0.2,0.15,9.6,0.35,
    'A. Evidence-Based Weight Derivation: Three Reasoning Chains',
    size=12,bold=True)

# Column headers
for x,lab in [(0.3,'Published Evidence'),(3.6,'Consequence for Biomarker'),(6.9,'Weight Adjustment')]:
    txt(sl,x,0.55,2.8,0.25,lab,size=9,bold=True,color=NAVY)

# ── Chain 1: TP53 → ctDNA ──
cy1=0.95; bw=2.6; bh=1.0
txt(sl,0.0,cy1+0.35,0.3,0.3,'(1)',size=12,bold=True,color=CORAL)
rrect(sl,0.3,cy1,bw,bh,LCORAL,CORAL)
txt(sl,0.3,cy1+0.1,bw,bh-0.2,'TP53 mutation\nimpairs apoptosis\n[11] Aubrey 2018',size=7,color=BODY)
arrow_r(sl,0.3+bw+0.05,cy1+bh/2,3.6-0.05,CORAL)
rrect(sl,3.6,cy1,bw,bh,LCORAL,CORAL)
txt(sl,3.6,cy1+0.1,bw,bh-0.2,'ctDNA release shifts to\nnecrotic + exosomal\npathways (different kinetics)',size=7,color=BODY)
arrow_r(sl,3.6+bw+0.05,cy1+bh/2,6.9-0.05,CORAL)
rrect(sl,6.9,cy1,bw,bh,CORAL,CORAL)
txt(sl,6.9,cy1+0.1,bw,bh-0.2,'ctDNA weight\n0.45 \u2192 0.30\n(\u2193 reliability)',size=8,bold=True,color=WHITE)

# ── Chain 2: Site → ALP ──
cy2=2.2
txt(sl,0.0,cy2+0.35,0.3,0.3,'(2)',size=12,bold=True,color=TEAL)
rrect(sl,0.3,cy2,bw,bh,LTEAL,TEAL)
txt(sl,0.3,cy2+0.1,bw,bh-0.2,'Tumor site determines\nbone ALP production\n[9,21] Lee 2004',size=7,color=BODY)
arrow_r(sl,0.3+bw+0.05,cy2+bh/2,3.6-0.05,TEAL)
rrect(sl,3.6,cy2,bw,bh,LTEAL,TEAL)
txt(sl,3.6,cy2+0.1,bw,bh-0.2,'Osseous: bone-remodeling\nALP present\nExtra-osseous: minimal ALP',size=7,color=BODY)
arrow_r(sl,3.6+bw+0.05,cy2+bh/2,6.9-0.05,TEAL)
rrect(sl,6.9,cy2,bw,bh,TEAL,TEAL)
txt(sl,6.9,cy2+0.1,bw,bh-0.2,'ALP weight\nOsseous: 0.25\nExtra-osseous: 0.10',size=8,bold=True,color=WHITE)

# ── Chain 3: STAG2 → ctDNA ──
cy3=3.45
txt(sl,0.0,cy3+0.35,0.3,0.3,'(3)',size=12,bold=True,color=PURPLE)
rrect(sl,0.3,cy3,bw,bh,LPURP,PURPLE)
txt(sl,0.3,cy3+0.1,bw,bh-0.2,'STAG2 mutation causes\nchromosomal instability\n[20] Solomon 2011',size=7,color=BODY)
arrow_r(sl,0.3+bw+0.05,cy3+bh/2,3.6-0.05,PURPLE)
rrect(sl,3.6,cy3,bw,bh,LPURP,PURPLE)
txt(sl,3.6,cy3+0.1,bw,bh-0.2,'Amplified constitutive\nctDNA shedding degrades\nclearance as response signal',size=7,color=BODY)
arrow_r(sl,3.6+bw+0.05,cy3+bh/2,6.9-0.05,PURPLE)
rrect(sl,6.9,cy3,bw,bh,PURPLE,PURPLE)
txt(sl,6.9,cy3+0.1,bw,bh-0.2,'STAG2 adjustment\nctDNA: \u22120.10\nLDH: +0.10',size=8,bold=True,color=WHITE)

# Footnote
txt(sl,0.3,4.6,9.4,0.25,
    'Each chain: published biological evidence \u2192 measurable consequence for biomarker interpretation \u2192 quantified weight adjustment',
    size=6.5,italic=True,color=GRAY)

# ── Panel B title ──
txt(sl,0.2,5.05,9.6,0.35,
    'B. Worked Example: How Reasoning Chains Produce Different BRS',
    size=12,bold=True)

# Shared context box
rrect(sl,1.0,5.5,8.0,0.6,GRAYBG,RGBColor(0xCC,0xCC,0xCC))
txt(sl,1.0,5.55,8.0,0.5,
    'Both patients: Osseous primary  \u2022  ctDNA cleared (s=1)  \u2022  LDH elevated (s=0)  \u2022  ALP normalized (s=1)',
    size=7.5,color=BODY)

# ── Patient A (left) ──
pax=0.5
rrect(sl,pax,6.3,4.3,0.5,LTEAL,TEAL)
txt(sl,pax,6.32,4.3,0.46,'Patient A: TP53-wild-type, Osseous',size=9,bold=True,color=BODY)
txt(sl,pax+0.2,6.95,3.0,0.25,'From Table 3 (TP53-wt, Osseous):',size=7,bold=True,color=NAVY,align=PP_ALIGN.LEFT)
txt(sl,pax+0.2,7.25,3.5,0.22,'ctDNA: w=0.45 \u00d7 s=1 = 0.45',size=7,color=BODY,align=PP_ALIGN.LEFT)
txt(sl,pax+0.2,7.50,3.5,0.22,'LDH:   w=0.30 \u00d7 s=0 = 0.00',size=7,color=BODY,align=PP_ALIGN.LEFT)
txt(sl,pax+0.2,7.75,3.5,0.22,'ALP:   w=0.25 \u00d7 s=1 = 0.25',size=7,color=BODY,align=PP_ALIGN.LEFT)
ln(sl,pax+0.2,8.02,pax+4.0,8.02,TEAL,1.0)
rrect(sl,pax+0.5,8.15,3.3,0.55,TEAL,TEAL,2.0)
txt(sl,pax+0.5,8.2,3.3,0.45,'BRS = 0.70',size=14,bold=True,color=WHITE)

# ── Patient B (right) ──
pbx=5.2
rrect(sl,pbx,6.3,4.3,0.5,LCORAL,CORAL)
txt(sl,pbx,6.32,4.3,0.46,'Patient B: TP53-mutant, Osseous',size=9,bold=True,color=BODY)
txt(sl,pbx+0.2,6.95,3.5,0.25,'From Table 3 (TP53-mut, Osseous):',size=7,bold=True,color=NAVY,align=PP_ALIGN.LEFT)
txt(sl,pbx+0.2,7.25,4.0,0.22,'ctDNA: w=0.30 \u00d7 s=1 = 0.30  \u2190 Chain (1)',size=7,color=BODY,align=PP_ALIGN.LEFT)
txt(sl,pbx+0.2,7.50,3.5,0.22,'LDH:   w=0.40 \u00d7 s=0 = 0.00',size=7,color=BODY,align=PP_ALIGN.LEFT)
txt(sl,pbx+0.2,7.75,3.5,0.22,'ALP:   w=0.30 \u00d7 s=1 = 0.30',size=7,color=BODY,align=PP_ALIGN.LEFT)
ln(sl,pbx+0.2,8.02,pbx+4.0,8.02,CORAL,1.0)
rrect(sl,pbx+0.5,8.15,3.3,0.55,CORAL,CORAL,2.0)
txt(sl,pbx+0.5,8.2,3.3,0.45,'BRS = 0.60',size=14,bold=True,color=WHITE)

# Delta annotation
ln(sl,pax+3.8,8.42,pbx+0.5,8.42,NAVY,2.0)
rrect(sl,4.2,8.25,1.6,0.35,WHITE,NAVY,1.5)
txt(sl,4.2,8.27,1.6,0.3,'\u0394 = 0.10',size=10,bold=True,color=NAVY)

# Bottom explanation
rrect(sl,1.0,8.9,8.0,0.95,RGBColor(0xFF,0xFF,0xF0),RGBColor(0xCC,0xCC,0xCC))
txt(sl,1.1,8.95,7.8,0.85,
    'Same biomarker pattern \u2192 different BRS.\n'
    'Chain (1) reduces ctDNA weight for TP53-mutant patients (impaired apoptosis \u2192 less reliable ctDNA clearance).\n'
    'The framework compensates by redistributing prognostic weight to LDH and ALP.',
    size=7,italic=True,color=BODY)

prs.save('./FigS1_Parameter_Justification.pptx')
print(f"FigS1_Parameter_Justification.pptx: {os.path.getsize('./FigS1_Parameter_Justification.pptx')/1024:.0f} KB")
if os.path.isdir('__pycache__'): shutil.rmtree('__pycache__')
